#!/usr/bin/env python3
"""
Build the "Making Time Uncertainty a First-Class Concept in Linux Timing"
40-minute presentation as a .pptx.

Fresh deck, written from scratch.

Usage:
    python3 build_deck.py [output.pptx]

Default output: time-uncertainty-error-bar.pptx (next to this script)
"""
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ----------------------------------------------------------------------------
# Theme
# ----------------------------------------------------------------------------
INK      = RGBColor(0x1A, 0x1D, 0x29)   # near-black text
PAPER    = RGBColor(0xFF, 0xFF, 0xFF)   # slide background
ACCENT   = RGBColor(0x2E, 0x6B, 0xE6)   # blue accent
ACCENT2  = RGBColor(0x14, 0xB8, 0xA6)   # teal
MUTED    = RGBColor(0x5B, 0x61, 0x72)   # muted gray text
LIGHT    = RGBColor(0xED, 0xF1, 0xFB)   # light panel
CODE_BG  = RGBColor(0x11, 0x16, 0x27)   # dark code panel
CODE_FG  = RGBColor(0xD6, 0xE2, 0xFF)   # code text
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

FONT      = "Calibri"
FONT_MONO = "Consolas"

SW = Inches(13.333)   # 16:9 widescreen
SH = Inches(7.5)

ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")


# ----------------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------------
def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.TOP, wrap=True, space_after=6):
    """runs: list of paragraphs; each paragraph is list of (text, size, bold,
    color, font, italic)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for (txt, size, bold, color, font, italic) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = font
            r.font.color.rgb = color
    return tb


def R(txt, size=18, bold=False, color=INK, font=FONT, italic=False):
    return (txt, size, bold, color, font, italic)


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _footer(slide, idx):
    _rect(slide, 0, SH - Inches(0.28), SW, Inches(0.28), ACCENT)
    _text(slide, Inches(0.5), SH - Inches(0.30), Inches(9), Inches(0.28),
          [[R("Time Uncertainty as a First-Class Concept", 10, False, WHITE)]],
          anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, SW - Inches(1.3), SH - Inches(0.30), Inches(0.8), Inches(0.28),
          [[R(str(idx), 10, True, WHITE)]], align=PP_ALIGN.RIGHT,
          anchor=MSO_ANCHOR.MIDDLE)


# ----------------------------------------------------------------------------
# Slide templates
# ----------------------------------------------------------------------------
def title_slide(prs, title, subtitle, presenter):
    s = _blank(prs)
    _set_bg(s, INK)
    _rect(s, 0, Inches(3.55), SW, Inches(0.09), ACCENT)
    _rect(s, 0, Inches(3.64), Inches(4.2), Inches(0.09), ACCENT2)
    _text(s, Inches(0.9), Inches(1.35), Inches(11.9), Inches(2.1),
          [[R(title, 38, True, WHITE)]])
    _text(s, Inches(0.92), Inches(3.85), Inches(11.5), Inches(0.8),
          [[R(subtitle, 22, False, RGBColor(0xC7, 0xD2, 0xF0))]])
    _text(s, Inches(0.92), Inches(5.9), Inches(11.5), Inches(1.0),
          [[R(presenter, 16, False, RGBColor(0x9A, 0xA6, 0xC4))]])
    return s


def section_slide(prs, part, title, idx):
    s = _blank(prs)
    _set_bg(s, ACCENT)
    _rect(s, 0, Inches(2.7), Inches(3.5), Inches(0.10), ACCENT2)
    _text(s, Inches(0.9), Inches(2.0), Inches(11), Inches(0.6),
          [[R(part.upper(), 20, True, RGBColor(0xC7, 0xD2, 0xF0))]])
    _text(s, Inches(0.88), Inches(2.95), Inches(11.5), Inches(2.0),
          [[R(title, 40, True, WHITE)]])
    _text(s, SW - Inches(1.3), SH - Inches(0.65), Inches(0.9), Inches(0.4),
          [[R(str(idx), 12, True, WHITE)]], align=PP_ALIGN.RIGHT)
    return s


def _header(s, title, kicker=None):
    _rect(s, 0, 0, Inches(0.18), Inches(1.25), ACCENT)
    if kicker:
        _text(s, Inches(0.55), Inches(0.28), Inches(12), Inches(0.4),
              [[R(kicker.upper(), 13, True, ACCENT)]])
        _text(s, Inches(0.55), Inches(0.62), Inches(12.2), Inches(0.9),
              [[R(title, 30, True, INK)]])
    else:
        _text(s, Inches(0.55), Inches(0.42), Inches(12.2), Inches(0.9),
              [[R(title, 30, True, INK)]])


def content_slide(prs, title, bullets, idx, kicker=None, notes=""):
    """bullets: list of (text, level) where level 0/1."""
    s = _blank(prs)
    _set_bg(s, PAPER)
    _header(s, title, kicker)
    paras = []
    for txt, lvl in bullets:
        if lvl == 0:
            paras.append([R("\u25B8  ", 18, True, ACCENT), R(txt, 18, False, INK)])
        else:
            paras.append([R("        \u2013  ", 16, False, MUTED),
                          R(txt, 16, False, MUTED)])
    _text(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(5.2), paras,
          space_after=10)
    _footer(s, idx)
    if notes:
        _notes(s, notes)
    return s


def statement_slide(prs, big, sub, idx, notes=""):
    s = _blank(prs)
    _set_bg(s, LIGHT)
    _rect(s, Inches(0.9), Inches(2.55), Inches(2.6), Inches(0.10), ACCENT)
    _text(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(2.0),
          [[R(big, 34, True, INK)]])
    if sub:
        _text(s, Inches(0.92), Inches(4.6), Inches(11.5), Inches(1.2),
              [[R(sub, 20, False, MUTED, FONT, True)]])
    _footer(s, idx)
    if notes:
        _notes(s, notes)
    return s


def code_slide(prs, title, code_lines, idx, kicker=None, caption=None, notes=""):
    s = _blank(prs)
    _set_bg(s, PAPER)
    _header(s, title, kicker)
    top = Inches(1.75)
    h = Inches(4.5) if caption else Inches(4.9)
    panel = _rect(s, Inches(0.7), top, Inches(11.9), h, CODE_BG)
    tf = panel.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = FONT_MONO
        r.font.size = Pt(15)
        r.font.color.rgb = CODE_FG
    if caption:
        _text(s, Inches(0.72), top + h + Inches(0.12), Inches(11.9), Inches(0.5),
              [[R(caption, 15, False, MUTED, FONT, True)]])
    _footer(s, idx)
    if notes:
        _notes(s, notes)
    return s


def table_slide(prs, title, headers, rows, idx, kicker=None, notes=""):
    s = _blank(prs)
    _set_bg(s, PAPER)
    _header(s, title, kicker)
    ncols = len(headers)
    nrows = len(rows) + 1
    left, top = Inches(0.7), Inches(1.8)
    avail = SH - top - Inches(0.5)          # keep clear of the footer bar
    row_h = min(Inches(0.5), int(avail / nrows))
    width, height = Inches(11.9), row_h * nrows
    gtbl = s.shapes.add_table(nrows, ncols, left, top, width, height).table
    for _r in gtbl.rows:
        _r.height = row_h
    body_sz = 13.5 if nrows <= 8 else 12
    # header
    for c, htext in enumerate(headers):
        cell = gtbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        cell.margin_left = Inches(0.12)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = htext
        r.font.bold = True; r.font.size = Pt(15); r.font.color.rgb = WHITE
        r.font.name = FONT
    # body
    for ri, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = gtbl.cell(ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PAPER if ri % 2 else LIGHT
            cell.margin_left = Inches(0.12)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            p = cell.text_frame.paragraphs[0]
            p.word_wrap = True
            r = p.add_run(); r.text = val
            r.font.size = Pt(body_sz)
            r.font.color.rgb = INK if c == 0 else MUTED
            r.font.bold = (c == 0)
            r.font.name = FONT
    _footer(s, idx)
    if notes:
        _notes(s, notes)
    return s


def two_col_slide(prs, title, left_head, left_items, right_head, right_items,
                  idx, kicker=None, notes=""):
    s = _blank(prs)
    _set_bg(s, PAPER)
    _header(s, title, kicker)
    top = Inches(1.85)
    colw = Inches(5.85)
    # left panel
    _rect(s, Inches(0.7), top, colw, Inches(4.7), LIGHT)
    _text(s, Inches(0.95), top + Inches(0.18), colw - Inches(0.5), Inches(0.5),
          [[R(left_head, 18, True, ACCENT)]])
    lp = [[R("\u2022  ", 15, True, ACCENT), R(t, 15, False, INK)] for t in left_items]
    _text(s, Inches(0.95), top + Inches(0.85), colw - Inches(0.5), Inches(3.6),
          lp, space_after=9)
    # right panel
    rx = Inches(0.7) + colw + Inches(0.4)
    _rect(s, rx, top, colw, Inches(4.7), RGBColor(0xEA, 0xF7, 0xF5))
    _text(s, rx + Inches(0.25), top + Inches(0.18), colw - Inches(0.5), Inches(0.5),
          [[R(right_head, 18, True, ACCENT2)]])
    rp = [[R("\u2022  ", 15, True, ACCENT2), R(t, 15, False, INK)] for t in right_items]
    _text(s, rx + Inches(0.25), top + Inches(0.85), colw - Inches(0.5), Inches(3.6),
          rp, space_after=9)
    _footer(s, idx)
    if notes:
        _notes(s, notes)
    return s


def image_slide(prs, image_path, idx, kicker=None, takeaway=None, notes=""):
    """Full-width figure slide. The image carries its own title; a short kicker
    sits on top and a one-line takeaway sits below."""
    s = _blank(prs)
    _set_bg(s, PAPER)
    if kicker:
        _rect(s, 0, Inches(0.18), Inches(0.18), Inches(0.55), ACCENT)
        _text(s, Inches(0.55), Inches(0.24), Inches(12), Inches(0.4),
              [[R(kicker.upper(), 13, True, ACCENT)]])
    img_w = Inches(11.7)
    img_h = Emu(int(img_w * 6.0 / 12.4))  # source figures are 12.4 x 6.0 in
    img_x = Emu(int((SW - img_w) / 2))
    img_y = Inches(0.82)
    s.shapes.add_picture(image_path, img_x, img_y, width=img_w, height=img_h)
    if takeaway:
        _text(s, Inches(0.7), img_y + img_h + Inches(0.06), Inches(11.9), Inches(0.5),
              [[R(takeaway, 15, True, INK)]], align=PP_ALIGN.CENTER)
    _footer(s, idx)
    if notes:
        _notes(s, notes)
    return s


# ----------------------------------------------------------------------------
# Build the deck
# ----------------------------------------------------------------------------
def build(path):
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    # 1 — Title
    title_slide(
        prs,
        "Making Time Uncertainty a First-Class Concept in Linux Timing",
        "Trustworthy timestamps for distributed and AI systems",
        "Time uncertainty as a first-class property  \u00b7  built on Linux + PTP",
    )

    # 2 — Hook
    content_slide(
        prs, "We make this decision constantly",
        [("event A at T\u2081   \u2192   event B at T\u2082   \u2192   \u201cA happened before B\u201d", 0),
         ("Ordering. Causality. \u201cWhich write won?\u201d \u201cWhat caused the stall?\u201d", 1),
         ("The hidden assumption: both timestamps are exact points.", 0),
         ("Clocks disagree. Timestamps are captured at different layers.", 1),
         ("Synchronization data goes stale between updates.", 1)],
        2, kicker="Opening",
        notes="Ask the room: who has debugged a distributed trace where events "
              "looked out of order? Almost every hand goes up. That is the whole talk.")

    # 3 — Two talks
    two_col_slide(
        prs, "Two talks, one story",
        "This talk", ["What is the error bar on a timestamp?",
                      "Timestamps as intervals",
                      "Explicit correctness criteria",
                      "\u201cCan I trust this time?\u201d"],
        "Prior talk \u2014 fast PHC access",
        ["How do we read the clock cheaply?",
         "PHC reads without syscall / PCIe overhead",
         "Userspace PHC approximation (Herm\u00f3\u00f0r)",
         "~5 \u00b5s syscall+PCIe read \u2192 ~60 ns"],
        3, kicker="Roadmap",
        notes="The prior talk made reading the PTP hardware clock cheap \u2014 the "
              "Herm\u00f3\u00f0r PoC approximates the device clock in userspace from the CPU "
              "counter, cutting a /dev/ptpX read from ~5 microseconds (syscall + "
              "PCIe) down to tens of nanoseconds. Note the shape: a daemon that "
              "cross-timestamps and estimates drift, a snapshot in shared memory, "
              "and a small client library \u2014 the same skeleton as ptp-uncertainty. "
              "This talk adds the missing half: the uncertainty of that time. "
              "Together they are the full story \u2014 fast access to the time and its "
              "associated error bar.")

    # ---- PART 1 ----
    section_slide(prs, "Part 1", "Precise is not the same as correct", 4)

    # 5
    content_slide(
        prs, "We spent two decades chasing precision",
        [("Sub-microsecond PTP is now common on data-center and AI fabrics", 0),
         ("Hardware timestamping in NICs, SmartNICs, DPUs", 0),
         ("Multi-node GPU traces spanning thousands of ranks", 0),
         ("The better our clocks got, the more we trusted them \u2014 implicitly", 0),
         ("Precision raised confidence faster than it raised correctness", 1)],
        5, kicker="The problem",
        notes="Precision is necessary but insufficient. Higher precision made "
              "silent errors both more likely to be relied on and harder to see.")

    # 6
    table_slide(
        prs, "Ways timestamps quietly mislead",
        ["Domain", "Naive claim", "What can go wrong"],
        [["Ordering", "A before B", "intervals overlap \u2192 order is unknown"],
         ["Causality", "A caused B", "off-by-one across nodes \u2192 false cause"],
         ["Compliance", "log is ordered", "cannot prove ordering is defensible"],
         ["Consistency", "commit A visible before B", "overlap breaks external consistency"],
         ["Profiling", "span A nested in span B", "cross-node spans misorder"],
         ["Leases / locks", "lease expired before regrant", "overlap \u2192 split brain / two leaders"],
         ["Telemetry", "counters sampled together", "inter-node skew \u2192 phantom correlations"],
         ["Forensics", "log A caused log B", "causal chain is indefensible"],
         ["Sensor fusion", "radar + camera aligned", "physical-world misalignment"],
         ["Finance", "trade timestamped correctly", "can't prove mandated window (MiFID II)"],
         ["Launch time load balancing", "packets will have their dedicated slots", "packet collisions"]],
        6, kicker="The problem",
        notes="Draft with the full menu of use cases \u2014 trim during review. The core "
              "three (ordering, causality, compliance) carry the argument; the rest "
              "map the blast radius: databases (external consistency / commit-wait), "
              "profiling (cross-node spans), coordination (leases \u2192 split brain), "
              "observability (telemetry skew), audit/forensics, cyber-physical "
              "(sensor fusion), and regulated timestamping (MiFID II). Auditors care "
              "about defensible ordering, not nanosecond vanity metrics.")

    # 7 — the problem, drawn: two overlapping error bars
    s = _blank(prs)
    _set_bg(s, LIGHT)
    _rect(s, Inches(0.9), Inches(1.02), Inches(2.6), Inches(0.10), ACCENT)
    _text(s, Inches(0.9), Inches(1.18), Inches(11.7), Inches(1.0),
          [[R("Two events 200 ns apart, each known to \u00b1150 ns.", 30, True, INK)]])

    AMBER_BG = RGBColor(0xFB, 0xE8, 0xC0)
    AMBER_DK = RGBColor(0xB4, 0x76, 0x0A)
    GUIDE = RGBColor(0xC2, 0xCB, 0xDD)
    A_FILL = ACCENT
    B_FILL = ACCENT2

    # time -> x mapping: t=0 ns at x0, 0.013 in per ns
    x0, per_ns = Inches(5.2), Inches(0.013)

    def X(ns):
        return int(x0 + int(round(ns * per_ns)))

    yA, yB, axis_y = Inches(3.35), Inches(4.35), Inches(5.05)
    bar_h, cap_h = Inches(0.18), Inches(0.36)

    # overlap strip (drawn first, behind everything)
    ox0, ox1 = X(50), X(150)
    _rect(s, ox0, Inches(2.95), ox1 - ox0, Inches(2.10), AMBER_BG)
    _text(s, ox0, Inches(3.0), ox1 - ox0, Inches(0.4),
          [[R("overlap", 14, True, AMBER_DK)]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # time axis
    _rect(s, X(-200), axis_y, X(400) - X(-200), Pt(2), MUTED)

    def error_bar(cy, lo, hi, c, fill, name):
        left, right = X(lo), X(hi)
        # guide from centre down to the axis
        _rect(s, X(c), cy, Pt(1.5), int(axis_y - cy), GUIDE)
        # interval band + end caps
        _rect(s, left, int(cy - bar_h / 2), right - left, bar_h, fill)
        _rect(s, left, int(cy - cap_h / 2), Pt(4), cap_h, fill)
        _rect(s, int(right - Pt(4)), int(cy - cap_h / 2), Pt(4), cap_h, fill)
        # centre marker (the reported timestamp)
        _rect(s, int(X(c) - Pt(3)), int(cy - cap_h / 2 - Inches(0.05)),
              Pt(6), int(cap_h + Inches(0.10)), INK)
        # event label to the left
        _text(s, int(left - Inches(0.62)), int(cy - Inches(0.2)),
              Inches(0.5), Inches(0.4), [[R(name, 18, True, fill)]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    error_bar(yA, -150, 150, 0, A_FILL, "A")
    error_bar(yB, 50, 350, 200, B_FILL, "B")

    # half-width dimension bracket under event A: 150 ns each side of the centre
    dim_y = int(yA + Inches(0.20))
    _rect(s, X(-150), dim_y, X(150) - X(-150), Pt(1.5), MUTED)
    for tx in (X(-150), X(0), X(150)):
        _rect(s, tx, int(dim_y - Inches(0.04)), Pt(1.5), Inches(0.11), MUTED)
    _text(s, X(-150), int(dim_y + Inches(0.05)), X(0) - X(-150), Inches(0.3),
          [[R("150 ns", 12, False, INK, FONT, True)]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _text(s, X(0), int(dim_y + Inches(0.05)), X(150) - X(0), Inches(0.3),
          [[R("150 ns", 12, False, INK, FONT, True)]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # gap annotation (200 ns) below the axis, between the two centres
    _text(s, X(0), int(axis_y + Inches(0.45)), X(200) - X(0), Inches(0.35),
          [[R("\u2190  200 ns apart  \u2192", 14, True, INK)]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _text(s, Inches(0.92), Inches(6.15), Inches(11.7), Inches(0.9),
          [[R("Their true order is unknown \u2014 and nothing in the timestamp told "
              "you so.", 20, False, MUTED, FONT, True)]])
    _footer(s, 7)
    _notes(s, "This is the aha. The numbers look confident; the reality is "
              "ambiguous. Each bar is the timestamp \u00b1 its uncertainty; because "
              "150 + 150 > 200, the intervals overlap (amber) and either event "
              "could have come first. We need the error bar to be part of the data.")

    # ---- PART 2 ----
    section_slide(prs, "Part 2", "A timestamp is an interval", 8)

    # 9
    content_slide(
        prs, "From a point to an interval",
        [("Instead of:  \u201cthe event happened at t\u201d", 0),
         ("Say:  \u201cthe event happened at t \u00b1 U\u201d", 0),
         ("U = the time-uncertainty bound at the moment of observation", 1),
         ("Also can be viewed as earliest-latest", 1),
         ("Now applications can ask explicit questions:", 0),
         ("Is ordering between A and B definite, or ambiguous?", 1),
         ("Is this window fresh enough for a compliance decision?", 1)],
        9, kicker="The model",
        notes="This is the core thesis of the abstract. Keep returning to it.")

    # 10 — ordering rule (drawn diagram, not ASCII art)
    s = _blank(prs)
    _set_bg(s, PAPER)
    _header(s, "The ordering rule", "The model")

    pnl_y, pnl_h = Inches(1.75), Inches(2.55)
    _rect(s, Inches(0.7), pnl_y, Inches(5.75), pnl_h, LIGHT)
    _rect(s, Inches(6.88), pnl_y, Inches(5.75), pnl_h, RGBColor(0xEA, 0xF7, 0xF5))
    _text(s, Inches(0.95), pnl_y + Inches(0.16), Inches(5.3), Inches(0.4),
          [[R("Point model", 16, True, MUTED)]])
    _text(s, Inches(7.13), pnl_y + Inches(0.16), Inches(5.3), Inches(0.4),
          [[R("Interval model", 16, True, ACCENT2)]])

    base_y = pnl_y + Inches(1.45)
    tick_h = Inches(0.56)
    # point model: baseline + single tick at t
    _rect(s, Inches(1.25), base_y, Inches(4.65), Pt(2.5), MUTED)
    _rect(s, Inches(3.55), int(base_y - tick_h / 2), Pt(3.5), tick_h, INK)
    _text(s, Inches(3.15), base_y + Inches(0.34), Inches(0.8), Inches(0.4),
          [[R("t", 16, True, INK)]], align=PP_ALIGN.CENTER)

    # interval model: baseline + bar [t-U, t+U] + centre tick
    _rect(s, Inches(7.4), base_y, Inches(4.85), Pt(2.5), MUTED)
    bar_x, bar_w = Inches(8.95), Inches(2.3)
    _rect(s, bar_x, int(base_y - Inches(0.12)), bar_w, Inches(0.24), ACCENT)
    _rect(s, bar_x, int(base_y - tick_h / 2), Pt(3.5), tick_h, ACCENT)
    _rect(s, int(bar_x + bar_w - Pt(3.5)), int(base_y - tick_h / 2), Pt(3.5), tick_h, ACCENT)
    _rect(s, int(bar_x + bar_w / 2), int(base_y - tick_h / 2 - Inches(0.05)),
          Pt(3.5), int(tick_h + Inches(0.1)), INK)
    # earliest / latest, above the interval end-caps
    _text(s, int(bar_x - Inches(0.7)), int(base_y - tick_h / 2 - Inches(0.34)),
          Inches(1.4), Inches(0.3), [[R("earliest", 13, False, ACCENT, FONT, True)]],
          align=PP_ALIGN.CENTER)
    _text(s, int(bar_x + bar_w - Inches(0.7)), int(base_y - tick_h / 2 - Inches(0.34)),
          Inches(1.4), Inches(0.3), [[R("latest", 13, False, ACCENT, FONT, True)]],
          align=PP_ALIGN.CENTER)
    _text(s, int(bar_x - Inches(0.55)), base_y + Inches(0.34), Inches(1.1), Inches(0.4),
          [[R("t\u2212U", 14, True, ACCENT)]], align=PP_ALIGN.CENTER)
    _text(s, int(bar_x + bar_w / 2 - Inches(0.4)), base_y + Inches(0.34), Inches(0.8), Inches(0.4),
          [[R("t", 14, True, INK)]], align=PP_ALIGN.CENTER)
    _text(s, int(bar_x + bar_w - Inches(0.55)), base_y + Inches(0.34), Inches(1.1), Inches(0.4),
          [[R("t+U", 14, True, ACCENT)]], align=PP_ALIGN.CENTER)

    # the rule band
    rule_y = Inches(4.6)
    _rect(s, Inches(0.7), rule_y, Inches(11.93), Inches(1.35), INK)
    _text(s, Inches(0.7), rule_y + Inches(0.22), Inches(11.93), Inches(0.6),
          [[R("A before B   only if   tA + UA  <  tB \u2212 UB", 26, True, WHITE, FONT_MONO)]],
          align=PP_ALIGN.CENTER)
    _text(s, Inches(0.7), rule_y + Inches(0.82), Inches(11.93), Inches(0.45),
          [[R("otherwise \u2192 ambiguous  \u00b7  unknown, not wrong", 18, False,
              RGBColor(0x8F, 0xE0, 0xD6), FONT, True)]],
          align=PP_ALIGN.CENTER)

    _text(s, Inches(0.72), Inches(6.12), Inches(11.9), Inches(0.5),
          [[R("The overlap region is where causality cannot be resolved from "
              "timestamps alone.", 15, False, MUTED, FONT, True)]])
    _footer(s, 10)
    _notes(s, "Draw this on a whiteboard if possible. The interval end-points have "
              "names: t\u2212U is the earliest the event could have happened, t+U the "
              "latest. 'Unknown' is a valid, useful answer a system can act on.")

    # 11 — is / is not
    two_col_slide(
        prs, "What uncertainty is \u2014 and is not",
        "U IS",
        ["A conservative bound on clock error, right now",
         "Derived from measurable sync state + configured limits",
         "An explicit input to application logic"],
        "U IS NOT",
        ["A statistical confidence interval (unless you add stats)",
         "A substitute for good clock design",
         "Zero just because PTP reports \u201cSLAVE\u201d"],
        11, kicker="The model",
        notes="Emphasize conservatism: we want a defensible upper bound, not a "
              "best guess.")

    # ---- PART 3 ----
    section_slide(prs, "Part 3", "Where does uncertainty come from?", 12)

    # 13 — the stack
    code_slide(
        prs, "Walk the stack, bottom to top",
        ["  6  Userspace observation delay   <- syscall + scheduling",
         "  5  Kernel -> userspace transfer  <- IRQ, softirq, copy",
         "  4  Timestamp capture point       <- HW vs SW",
         "  3  Network                       <- queueing, routing, hop count",
         "  2  PHC clock                     <- offset, clock resolution",
         "  1  Oscillator                    <- drift, temperature"],
        13, kicker="Anatomy",
        caption="Each layer adds error or delay. Only some layers expose metrics.",
        notes="Walk slowly. The point: 'the event time' means different things at "
              "different layers even when everything is 'synced'.")

    # 14 — the full budget of sources
    table_slide(
        prs, "The uncertainty budget: every source",
        ["Source of error", "What it contributes", "In the bound?"],
        [["Grandmaster class",
          "advertised GM quality (clockClass / accuracy)",
          "error bound we are told"],
         ["Offset from grandmaster",
          "our offset, measured from the t1\u2013t4 exchange",
          "yes \u2014 |offset| term"],
         ["Local oscillator drift",
          "free-run drift between syncs (ppb\u2013ppm)",
          "yes \u2014 drift term"],
         ["Holdover",
          "drift with no discipline (ptp4l down)",
          "yes \u2014 anchor held, drift grows"],
         ["Clock read delays",
          "latency to read the PHC / clock_gettime",
          "partial \u2014 capture-point dependent"],
         ["Clock resolution",
          "granularity of the timestamp tick",
          "bounded by tick size"],
         ["Accumulated path-element delays",
          "residence + queueing over each hop",
          "compensated by PTP (mean)"],
         ["Link asymmetry",
          "up and down paths differ",
          "unmodeled \u2014 conservative gap"]],
        14, kicker="Anatomy",
        notes="A consolidated budget before we zoom into individual layers. Two "
              "sources come from the source clock itself, and they differ in kind: "
              "the grandmaster's class (clockClass / clockAccuracy) is an advertised "
              "bound on the source quality, while the offset from the grandmaster is "
              "our own value, measured locally from the t1-t4 timestamp exchange \u2014 "
              "that measured |offset| is what enters the bound. (stepsRemoved is "
              "neither: just the hop count to the GM.) Two come "
              "from our local oscillator: its free-run drift between sync messages, "
              "and holdover when discipline is lost entirely. Two are properties of "
              "reading the clock: the access latency of a PHC or clock_gettime "
              "read, and the finite resolution of the tick itself. And two live in "
              "the network: the delays accumulated across every transparent/"
              "boundary clock and queue on the path (whose mean PTP measures and "
              "compensates), and link asymmetry \u2014 the up/down difference PTP cannot "
              "see, which stays an unmodeled conservative gap. The rest of Part 3 "
              "drills into the ones that dominate.")

    # 15 — oscillator & PHC clock
    content_slide(
        prs, "Layers 1\u20132: oscillator & PHC clock",
        [("Oscillator (layer 1): TCXO/OCXO stability (ppb), temperature sensitivity, free-run drift", 0),
         ("PHC clock (layer 2): the disciplined hardware clock we actually read", 0),
         ("offsetFromMaster \u2014 offset estimate", 1),
         ("clock resolution \u2014 granularity of a single tick", 1),
         ("PHC read latency \u2014 the cost of reading it", 1),
         ("The servo converges toward zero offset; it never mathematically reaches it", 0)],
        15, kicker="Anatomy",
        notes="Layer 1 is the oscillator alone \u2014 stability in ppb and temperature "
              "sensitivity, free-running drift. Layer 2 is the PHC clock it feeds: "
              "offsetFromMaster is an estimate (not a truth), and the clock has a "
              "finite resolution (tick granularity) plus a read latency. PTP "
              "believes something about your clock; that belief is an input to U, "
              "not a guarantee.")

    # 16 — capture point table
    table_slide(
        prs, "Layer 4: where was the timestamp taken?",
        ["Method", "Typical error", "Notes"],
        [["PHC hardware RX timestamp", "smallest", "needs NIC/driver support"],
         ["SO_TIMESTAMPING (software)", "larger", "IRQ + softirq delay"],
         ["clock_gettime() in app", "largest", "syscall + scheduling jitter"]],
        16, kicker="Anatomy",
        notes="Critical for profiling: GPU, NIC and CPU may live in different "
              "timestamp domains even when 'synced'.")

    # 17 — network & kernel path
    content_slide(
        prs, "Layers 3 & 5: network and kernel path",
        [("Network (layer 3): each hop adds queueing + routing delay", 0),
         ("PTP compensates the mean path delay \u2014 so it is not in the bound", 1),
         ("Residual is link asymmetry (up/down differ), accumulating with hop count \u2014 unmodeled", 1),
         ("Kernel path (layer 5): Linux exposes useful data \u2014 but not as one package:", 0),
         ("PTP_SYS_OFFSET / _EXTENDED \u2014 PHC vs system clock", 1),
         ("SO_TIMESTAMPING \u2014 ingress / egress timestamps", 1),
         ("ptp4l management API \u2014 offset, delay, ingress time", 1)],
        17, kicker="Anatomy",
        notes="Layer 3 is the network: every hop \u2014 transparent/boundary clocks and "
              "queues \u2014 adds delay. PTP measures and compensates the mean path "
              "delay, so it is not in the bound; only the residual up/down asymmetry "
              "is, it grows with hop count, and we do not get a clean measurement of "
              "it. Layer 5 is the kernel path: the ingredients exist (PTP_SYS_OFFSET, "
              "SO_TIMESTAMPING, ptp4l mgmt), but the single 'current uncertainty' "
              "API does not. Foreshadow the kernel-gaps section.")

    # 18 — staleness (the big one)
    code_slide(
        prs, "Layer that everyone forgets: staleness",
        ["  Between PTP updates, the clock drifts.",
         "",
         "  U_drift  =  elapsed_time  x  max_drift_rate",
         "",
         "  Example:  100 ppm  x  1 s   =   100 microseconds",
         "",
         "  A perfect offset at sync ingress still becomes",
         "  a GROWING interval a moment later."],
        18, kicker="Anatomy",
        caption="Staleness turns a point sync measurement into a growing interval "
                "\u2014 this is what the model daemon implements.",
        notes="Many tools show offset but never show the interval growing since "
              "the last correction. This is the term the daemon makes explicit.")

    # ---- PART 4 ----
    section_slide(prs, "Part 4", "Turning PTP state into a bound", 19)

    # 20 — datasets
    table_slide(
        prs, "PTP answers: what does the protocol believe?",
        ["Dataset (via ptp4l)", "Provides"],
        [["CURRENT_DATA_SET", "offsetFromMaster, meanPathDelay, stepsRemoved"],
         ["PORT_DATA_SET", "port state, logSyncInterval"],
         ["TIME_STATUS_NP", "sync ingress time, grandmaster identity"],
         ["PORT_HWCLOCK_NP", "PHC index (optional)"]],
        20, kicker="PTP \u2192 bound",
        notes="PTP does NOT answer 'what is my application-level timestamp error "
              "bar?'. We build that on top.")

    # 21 — ingress anchor
    code_slide(
        prs, "Ingress time is the drift anchor",
        ["  TIME_STATUS_NP.ingress_time",
         "     = when the last sync event arrived (PHC time)",
         "",
         "  drift_ns = (time_now - ingress_time)",
         "                 * max_drift_ppb / 1e9",
         "",
         "  Uncertainty grows between sync messages",
         "  even when the offset display looks flat."],
        21, kicker="PTP \u2192 bound",
        notes="This makes 'staleness' concrete: anchor at ingress, extrapolate at "
              "read time using a worst-case drift bound.")

    # 22 — the formula
    statement_slide(
        prs,
        "total_uncertainty_ns = |offset| + drift + clock_resolution + capture_point_error",
        "Example: 50 ns drift + 8 ns resolution + 40 ns capture \u2248 98 ns bound.",
        22,
        notes="Three residual terms after the offset correction is applied: drift "
              "(staleness), clock resolution (tick granularity), and capture-point "
              "error (HW vs SW stamp). Mean path delay is measured and compensated "
              "by PTP inside the offset, so it is not added; the offset itself is a "
              "correction, not an uncertainty. Residual path asymmetry is real but "
              "unmodeled here \u2014 a known conservative gap. A practical model, not a "
              "full Allan-deviation analysis.")

    # 23 — port state trust
    table_slide(
        prs, "Sync status is a prerequisite, not a proof",
        ["Port state", "Meaning for uncertainty"],
        [["SLAVE", "actively disciplining \u2014 bounds meaningful"],
         ["UNCALIBRATED", "converging \u2014 bounds unstable"],
         ["LISTENING / FAULTY", "do not trust ordering claims"],
         ["ptp4l disconnected", "hold last anchor + drift keeps growing"]],
        23, kicker="PTP \u2192 bound",
        notes="'Synchronized' is necessary but not sufficient. The last row is the "
              "interesting failure mode the daemon handles explicitly.")

    # ---- PART 5 ----
    section_slide(prs, "Part 5", "A model implementation: ptp-uncertainty", 24)

    # 25 — architecture
    code_slide(
        prs, "Architecture",
        ["  ptp4l",
         "    |  Unix management socket (poll offset/delay/ingress)",
         "    v",
         "  ptp_unc_dmn        <- daemon: collects + anchors",
         "    |  /ptp_uncertainty (POSIX shared memory)",
         "    v",
         "  libptp_unc.so      <- client lib: extrapolates at read time",
         "    |",
         "    v",
         "  applications       <- get a live bound, no PTP code needed"],
        25, kicker="Implementation",
        caption="No direct PHC access required. A live bound, not just the last "
                "offset snapshot.",
        notes="Transition from theory to something the audience can actually run.")

    # 26 — what daemon collects
    table_slide(
        prs, "What the daemon collects",
        ["Field", "Source"],
        [["offset, delay, steps removed", "CURRENT_DATA_SET"],
         ["port state, sync interval", "PORT_DATA_SET"],
         ["ingress time, grandmaster id", "TIME_STATUS_NP"],
         ["PHC index (optional)", "PORT_HWCLOCK_NP"]],
        26, kicker="Implementation",
        notes="Auto features: poll interval derived from logSyncInterval (>=2x "
              "sync rate); optional PHC index autodetection.")

    # 27 — client API
    code_slide(
        prs, "Client API",
        ["  struct ptp_unc_handle *h = ptp_unc_open();",
         "  struct ptp_uncertainty u;",
         "",
         "  ptp_unc_get(h, &u);            // uncertainty at now",
         "  ptp_unc_get_at(h, &u, t_mono); // at a given monotonic t",
         "",
         "  // u.total_uncertainty_ns, u.drift_ns",
         "  // u.is_synchronized, u.ptp4l_connected",
         "  // u.age_ns  (snapshot freshness != drift anchor)",
         "",
         "  ptp_unc_close(h);"],
        27, kicker="Implementation",
        caption="Clients extrapolate at read time; age_ns (snapshot freshness) is "
                "distinct from the drift term.",
        notes="Clarify snapshot age vs drift \u2014 a common confusion and a good "
              "teaching moment.")

    # 28 — failure behavior
    code_slide(
        prs, "Behavior under failure = explicit correctness",
        ["  // ptp4l restarts or disconnects:",
         "  //   - preserve last valid sync anchor",
         "  //   - set ptp4l_connected = 0",
         "  //   - drift KEEPS growing from the held anchor",
         "",
         "  if (!u.ptp4l_connected ||",
         "      u.total_uncertainty_ns > threshold)",
         "      mark_event_ordering_untrusted();"],
        28, kicker="Implementation",
        caption="The application decides what 'trustworthy enough' means \u2014 the "
                "daemon just keeps the bound honest.",
        notes="This is the abstract's phrase 'evaluate temporal correctness "
              "explicitly' turned into three lines of code.")

    # 29 — industry validation: Meta fbclock
    two_col_slide(
        prs, "You don\u2019t have to take my word for it: Meta\u2019s fbclock",
        "Meta fbclock  (deployed at scale, 2022)",
        ["Returns a Window of Uncertainty: [earliest_ns, latest_ns]",
         "Go daemon reads ptp4l + PHC \u2192 shared memory \u2192 C library",
         "Statistical k\u00b7\u03c3 window (Meta targets 6-nines certainty)",
         "Holdover: drift from PHC freq-adjust history + temp/vibration telemetry"],
        "ptp-uncertainty  (this talk)",
        ["Returns total_uncertainty_ns \u2192 [t\u2212U, t+U]",
         "ptp_unc_dmn reads ptp4l \u2192 /ptp_uncertainty SHM \u2192 libptp_unc.so",
         "Conservative worst-case bound (layer stats on top if wanted)",
         "Holdover: hold last anchor, drift grows at max_drift_ppb"],
        29, kicker="Industry validation",
        notes="The interval model is not academic. Meta's fbclock ships 'Window of "
              "Uncertainty' [earliest, latest] to a fleet of PTP clients using the "
              "same daemon+shmem+library shape this talk arrives at independently. "
              "Two honest differences: (1) Meta computes a statistical sigma-based "
              "window (M = clockaccuracy + |offset| + stddev(offset); W = mean(M) + "
              "k*stddev(M), k->4), we default to a conservative worst-case bound; "
              "(2) Meta estimates holdover drift empirically from PHC frequency-"
              "adjustment history (1.5 * mean(|freqchange|)) plus temperature and "
              "vibration telemetry, we use a configured max_drift_ppb. Same skeleton, "
              "different knobs. Source: 'How Precision Time Protocol is being deployed "
              "at Meta', Meta Engineering, 2022; github.com/facebook/time/fbclock.")

    # ---- PART 6 ----
    section_slide(prs, "Part 6", "Measured on real oscillators", 30)

    # 30 — staleness in action (sawtooth)
    image_slide(
        prs, os.path.join(ASSETS, "fig_sawtooth.png"), 31, kicker="Measured",
        takeaway="A perfect sync becomes a growing interval within one second \u2014 "
                 "this is the drift term, captured live.",
        notes="Real capture from watch_uncertainty against a live ptp4l at 1 Hz "
              "sync, 10 ppm drift bound. Mean path delay is compensated by PTP, so "
              "the bound is |offset| + drift; the dashed line is the tiny |offset| "
              "floor (tens of ns) and the sawtooth is staleness accumulating "
              "between sync messages and resetting at each ingress. This is slide "
              "17's claim, measured.")

    # 31 — oscillator comparison
    image_slide(
        prs, os.path.join(ASSETS, "fig_compare.png"), 32, kicker="Measured",
        takeaway="Same PTP quality, four oscillator classes \u2014 the drift bound "
                 "alone spans two orders of magnitude.",
        notes="Every run has the same tiny |offset| floor (tens of ns) because path "
              "delay is compensated by PTP. Only max_drift_ppb changes: 100 ppb "
              "(OCXO), 1 ppm (TCXO), 10 ppm, 100 ppm. Log scale. The envelope is "
              "set by staleness, not by PTP.")

    # 32 — peak budget bars
    image_slide(
        prs, os.path.join(ASSETS, "fig_budget.png"), 33, kicker="Measured",
        takeaway="Between sync messages, oscillator holdover \u2014 not PTP \u2014 "
                 "decides whether your bound is 0.3 \u00b5s or 120 \u00b5s.",
        notes="Peak-moment budget per class, as |offset| + drift. For the OCXO, "
              "offset and drift are comparable (~0.3 us total); for the basic XO, "
              "drift is essentially 100%. Same protocol, same network, the peak "
              "bound runs from ~0.3 us to ~120 us from the crystal alone. This is "
              "why oscillator stability must be part of the model and why the "
              "kernel not exposing it (next section) is the real gap.")

    # ---- PART 7 ----
    section_slide(prs, "Part 7", "The full story: fast time + its uncertainty", 34)

    # 35 — bridge
    statement_slide(
        prs,
        "Fast access gives you the time cheaply. Uncertainty tells you how far to trust it.",
        "The prior talk made the read fast; this one puts an error bar on it \u2014 fast time and its uncertainty, together.",
        35,
        notes="Explicit bridge to the prior fast-PHC-access talk (Herm\u00f3\u00f0r). Both "
              "halves share a daemon + shared-memory + client-library shape and "
              "compose into one primitive: a timestamp you can read cheaply \u2014 tens "
              "of nanoseconds instead of microseconds \u2014 and trust explicitly.")

    # 36 — the two halves compose
    table_slide(
        prs, "Fast time + its uncertainty compose",
        ["Capability", "Fast read alone", "+ uncertainty"],
        [["Low-overhead timestamp", "yes", "yes"],
         ["Compare two events", "point compare", "interval compare"],
         ["Know when order is ambiguous", "no", "explicit"],
         ["Gate decisions on trust", "no", "threshold on U"],
         ["Merge / correlation window", "fixed", "adapts to live U"]],
        36, kicker="The full story",
        notes="Fast reads (prior talk) and the uncertainty bound (this talk) are "
              "two halves of one time primitive: read the clock cheaply, and know "
              "how wrong it might be. Example: two events 300 ns apart, each U ~ "
              "500 ns \u2014 the fast read alone says 'A before B'; with U the order is "
              "honestly ambiguous.")

    # ---- PART 8 ----
    section_slide(prs, "Part 8", "Linux APIs: what we have, what we lack", 37)

    # 32 — have vs missing
    two_col_slide(
        prs, "The kernel gives ingredients, not the recipe",
        "Available today",
        ["SO_TIMESTAMPING \u2014 HW/SW RX/TX timestamps",
         "PTP_SYS_OFFSET(_EXTENDED) \u2014 PHC \u2194 system",
         "ptp4l mgmt socket \u2014 offset, delay, ingress",
         "/dev/ptpN \u2014 freq adjust, ext timestamps"],
        "Not exposed / not standardized",
        ["Oscillator stability (Allan deviation)",
         "Calibration age + temperature model",
         "Per-read PHC access-latency bound",
         "A unified \u201cstaleness since last discipline\u201d"],
        38, kicker="Kernel limits",
        notes="Matches the abstract's conclusion: frequency offset is available, "
              "but oscillator stability and calibration staleness are not \u2014 so "
              "precise bounds still need operator config and hardware knowledge. "
              "The measurements just proved how much that missing oscillator term "
              "matters. And Meta's fbclock is the existence proof: to run holdover "
              "in production they estimate drift from PHC frequency-adjustment "
              "history and calibrate it with temperature and vibration telemetry "
              "\u2014 exactly the data the kernel does not hand you. Do not oversell "
              "automation.")

    # ---- CLOSE ----
    content_slide(
        prs, "Takeaways",
        [("Timestamps are intervals \u2014 design systems around t \u00b1 U", 0),
         ("PTP provides inputs, not a complete uncertainty model", 0),
         ("Staleness / drift is the term most tools ignore \u2014 don\u2019t", 0),
         ("Explicit bounds enable explicit correctness for ordering & compliance", 0),
         ("The full story: fast access to the time and its associated uncertainty", 0)],
        39, kicker="Closing",
        notes="End on the bridge sentence: better profiles need better time \u2014 "
              "and better time needs error bars.")

    # 34 — thanks / QA
    s = _blank(prs)
    _set_bg(s, INK)
    _rect(s, 0, Inches(3.0), SW, Inches(0.09), ACCENT)
    _rect(s, 0, Inches(3.09), Inches(4.2), Inches(0.09), ACCENT2)
    _text(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.0),
          [[R("Questions?", 44, True, WHITE)]])
    _text(s, Inches(0.92), Inches(3.4), Inches(11.5), Inches(1.4),
          [[R("Treat time as an interval. Make the error bar part of the data.",
              20, False, RGBColor(0xC7, 0xD2, 0xF0))]])
    _text(s, Inches(0.92), Inches(5.6), Inches(11.5), Inches(0.8),
          [[R("ptp-uncertainty  \u00b7  daemon + libptp_unc.so + watch_uncertainty",
              16, False, RGBColor(0x9A, 0xA6, 0xC4), FONT_MONO)]])
    _notes(s, "Offer to collaborate with profiling-tool authors. Have the demo "
              "plot ready as a backup.")

    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)
    return len(prs.slides._sldIdLst)


if __name__ == "__main__":
    here = os.path.dirname(os.path.abspath(__file__))
    default = os.path.join(here, "time-uncertainty-error-bar.pptx")
    out = sys.argv[1] if len(sys.argv) > 1 else default
    n = build(out)
    print(f"Wrote {n} slides -> {out}")
